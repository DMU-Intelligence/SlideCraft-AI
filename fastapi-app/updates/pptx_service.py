from __future__ import annotations

"""PPTX 생성 서비스 — element 기반

ProjectState.slides 의 PageLayout.elements 를 순서대로 렌더링합니다.
지원 element: text_box / shape / bullet_list
"""

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models.project_state import ProjectState
from ..schemas.generate import (
    BulletListElement,
    PageLayout,
    ShapeElement,
    SlideContent,
    TextBoxElement,
    SlideElement,
)

# ── 슬라이드 크기 (16:9) ───────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

_ALIGN_MAP: dict[str, Any] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


# ── 헬퍼 ───────────────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "FFFFFF"  # fallback
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _set_bg(slide: Any, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


# ── Element 렌더러 ─────────────────────────────────────────────────────────────

def _render_text_box(slide: Any, elem: TextBoxElement) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(elem.left), Inches(elem.top),
        Inches(elem.width), Inches(elem.height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = _ALIGN_MAP.get(elem.align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = elem.text
    run.font.name = elem.font_name
    run.font.size = Pt(elem.font_size)
    run.font.bold = elem.font_bold
    run.font.color.rgb = _hex_to_rgb(elem.font_color)


def _render_shape(slide: Any, elem: ShapeElement) -> None:
    shape_map = {
        "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "round_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    }
    shape = slide.shapes.add_shape(
        shape_map.get(elem.shape_type, MSO_AUTO_SHAPE_TYPE.RECTANGLE),
        Inches(elem.left), Inches(elem.top),
        Inches(elem.width), Inches(elem.height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(elem.fill_color)
    shape.line.fill.background()


def _render_bullet_list(slide: Any, elem: BulletListElement) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(elem.left), Inches(elem.top),
        Inches(elem.width), Inches(elem.height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    bullet_rgb = _hex_to_rgb(elem.bullet_color)
    font_rgb = _hex_to_rgb(elem.font_color)

    for i, item_text in enumerate(elem.items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        # 불릿 기호
        dot = p.add_run()
        dot.text = f"{elem.bullet_char}  "
        dot.font.name = elem.font_name
        dot.font.size = Pt(elem.font_size)
        dot.font.color.rgb = bullet_rgb
        dot.font.bold = True

        # 본문
        body = p.add_run()
        body.text = item_text
        body.font.name = elem.font_name
        body.font.size = Pt(elem.font_size)
        body.font.color.rgb = font_rgb
        body.font.bold = False


def _render_element(slide: Any, element: SlideElement) -> None:
    if isinstance(element, TextBoxElement):
        _render_text_box(slide, element)
    elif isinstance(element, ShapeElement):
        _render_shape(slide, element)
    elif isinstance(element, BulletListElement):
        _render_bullet_list(slide, element)


# ── 페이지(슬라이드) 빌더 ──────────────────────────────────────────────────────

def _build_page(prs: Presentation, page: PageLayout) -> None:
    slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, page.background)
    for element in page.elements:
        _render_element(slide, element)


# ── 공개 인터페이스 ────────────────────────────────────────────────────────────

class PptxGenerator:
    """ProjectState.slides → .pptx bytes"""

    def generate(self, state: ProjectState) -> bytes:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        for slide_content in state.slides:
            for page in slide_content.pages:
                _build_page(prs, page)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

from __future__ import annotations

import unittest
from types import SimpleNamespace

from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.schemas.generate import LineElement
from app.services.template_extractor import RawBounds, TemplateExtractor, _IDENTITY_TRANSFORM


class _FakeFill:
    def __init__(self, fill_type: object | None, rgb: str | None = None) -> None:
        self.type = fill_type
        self.fore_color = SimpleNamespace(rgb=rgb)


class _FakeLine:
    def __init__(self, rgb: str | None = None, width: int = 0) -> None:
        self.fill = _FakeFill(MSO_FILL.SOLID if rgb else None, rgb)
        self.width = width


class _FakePoint:
    def __init__(self, x: float = 0, y: float = 0) -> None:
        self.x = x
        self.y = y


class _FakeSize:
    def __init__(self, cx: float = 0, cy: float = 0) -> None:
        self.cx = cx
        self.cy = cy


class _FakeXfrm:
    def __init__(
        self,
        *,
        x: float,
        y: float,
        cx: float,
        cy: float,
        ch_off_x: float,
        ch_off_y: float,
        ch_ext_cx: float,
        ch_ext_cy: float,
    ) -> None:
        self.x = x
        self.y = y
        self.cx = cx
        self.cy = cy
        self.chOff = _FakePoint(ch_off_x, ch_off_y)
        self.chExt = _FakeSize(ch_ext_cx, ch_ext_cy)


class _FakeShape:
    def __init__(
        self,
        *,
        shape_type: object,
        left: float = 0,
        top: float = 0,
        width: float = 0,
        height: float = 0,
        fill: object | None = None,
        line: object | None = None,
        has_text_frame: bool = False,
        text_frame: object | None = None,
        auto_shape_type: object | None = None,
        shapes: list[object] | None = None,
        xfrm: object | None = None,
        begin_x: float | None = None,
        begin_y: float | None = None,
        end_x: float | None = None,
        end_y: float | None = None,
        name: str = "shape",
    ) -> None:
        self.shape_type = shape_type
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.fill = fill
        self.line = line or _FakeLine()
        self.has_text_frame = has_text_frame
        self.text_frame = text_frame
        self.auto_shape_type = auto_shape_type
        self.shapes = shapes or []
        self._element = SimpleNamespace(xfrm=xfrm)
        self.name = name
        if begin_x is not None:
            self.begin_x = begin_x
        if begin_y is not None:
            self.begin_y = begin_y
        if end_x is not None:
            self.end_x = end_x
        if end_y is not None:
            self.end_y = end_y


class TemplateExtractorTests(unittest.TestCase):
    def setUp(self) -> None:
        self.extractor = TemplateExtractor()

    def test_normalize_box_bounds_uses_source_slide_ratio(self) -> None:
        bounds = self.extractor._normalize_box_bounds(
            RawBounds(x=2, y=1, w=4, h=2),
            source_width=20,
            source_height=10,
        )

        self.assertEqual(
            bounds,
            {
                "x": 1.33,
                "y": 0.75,
                "w": 2.67,
                "h": 1.5,
            },
        )

    def test_group_shape_recurses_and_applies_transform(self) -> None:
        child = _FakeShape(
            shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE,
            left=10,
            top=10,
            width=20,
            height=10,
            fill=_FakeFill(MSO_FILL.SOLID, "FF0000"),
            name="child",
        )
        group = _FakeShape(
            shape_type=MSO_SHAPE_TYPE.GROUP,
            shapes=[child],
            xfrm=_FakeXfrm(
                x=50,
                y=30,
                cx=200,
                cy=100,
                ch_off_x=0,
                ch_off_y=0,
                ch_ext_cx=100,
                ch_ext_cy=50,
            ),
            name="group",
        )

        elements = self.extractor._extract_shape_element(
            group,
            source_width=200,
            source_height=100,
            transform=_IDENTITY_TRANSFORM,
            slide_debug={"extracted": [], "skipped": []},
            source_path="slide[1].shape[1]",
        )

        self.assertEqual(len(elements), 1)
        self.assertEqual(elements[0]["type"], "shape")
        self.assertEqual(elements[0]["x"], 4.67)
        self.assertEqual(elements[0]["y"], 3.75)
        self.assertEqual(elements[0]["w"], 2.67)
        self.assertEqual(elements[0]["h"], 1.5)

    def test_freeform_thin_shape_is_normalized_to_line(self) -> None:
        freeform = _FakeShape(
            shape_type=MSO_SHAPE_TYPE.FREEFORM,
            left=20,
            top=20,
            width=200,
            height=2,
            fill=_FakeFill(MSO_FILL.SOLID, "00FF00"),
            line=_FakeLine("00FF00", width=12700),
        )

        elements = self.extractor._extract_freeform_shape(
            freeform,
            source_width=400,
            source_height=200,
            transform=_IDENTITY_TRANSFORM,
            slide_debug={"extracted": [], "skipped": []},
            source_path="slide[1].shape[2]",
        )

        self.assertEqual(len(elements), 1)
        self.assertEqual(elements[0]["type"], "line")
        self.assertEqual(elements[0]["h"], 0.0)
        self.assertGreater(elements[0]["w"], 0.0)

    def test_connector_line_bounds_preserve_zero_height(self) -> None:
        line_shape = _FakeShape(
            shape_type=MSO_SHAPE_TYPE.LINE,
            line=_FakeLine("112233", width=12700),
            begin_x=20,
            begin_y=40,
            end_x=80,
            end_y=40,
        )

        element = self.extractor._extract_line_shape(
            line_shape,
            source_width=100,
            source_height=100,
            transform=_IDENTITY_TRANSFORM,
        )

        self.assertIsNotNone(element)
        assert element is not None
        self.assertEqual(element["type"], "line")
        self.assertEqual(element["h"], 0.0)
        self.assertEqual(element["w"], 8.0)

    def test_zero_dimension_auto_shape_is_treated_as_line(self) -> None:
        line_shape = _FakeShape(
            shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE,
            left=10,
            top=20,
            width=50,
            height=0,
            line=_FakeLine("445566", width=12700),
        )

        elements = self.extractor._extract_auto_shape(
            line_shape,
            source_width=100,
            source_height=100,
            transform=_IDENTITY_TRANSFORM,
            slide_debug={"extracted": [], "skipped": []},
            source_path="slide[1].shape[3]",
        )

        self.assertEqual(len(elements), 1)
        self.assertEqual(elements[0]["type"], "line")
        self.assertEqual(elements[0]["h"], 0.0)


class LineElementSchemaTests(unittest.TestCase):
    def test_line_element_accepts_zero_height(self) -> None:
        line = LineElement(x=1.0, y=1.0, w=10.0, h=0.0)
        self.assertEqual(line.h, 0.0)

    def test_line_element_rejects_zero_vector(self) -> None:
        with self.assertRaises(ValueError):
            LineElement(x=1.0, y=1.0, w=0.0, h=0.0)


if __name__ == "__main__":
    unittest.main()

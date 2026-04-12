from __future__ import annotations

from dataclasses import dataclass
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

from ..core.logging import get_logger
from ..schemas.generate import PageLayout, SLIDE_HEIGHT_INCHES, SLIDE_WIDTH_INCHES

_LOGGER = get_logger(__name__)

_DEFAULT_BACKGROUND = "#FFFFFF"
_DEFAULT_FONT_NAME = "Malgun Gothic"
_DEFAULT_FONT_SIZE = 16
_DEFAULT_FONT_COLOR = "#000000"
_DEFAULT_BULLET_COLOR = "#2563EB"
_DEFAULT_LINE_COLOR = "#D9D9D9"
_DEFAULT_LINE_WIDTH = 1.0
_FREEFORM_LINE_THICKNESS_MAX = 0.14
_FREEFORM_LINE_ASPECT_RATIO = 7.0


@dataclass(frozen=True)
class AxisTransform:
    scale: float = 1.0
    translate: float = 0.0

    def apply(self, value: float) -> float:
        return (value * self.scale) + self.translate

    def compose(self, other: "AxisTransform") -> "AxisTransform":
        return AxisTransform(
            scale=self.scale * other.scale,
            translate=(self.scale * other.translate) + self.translate,
        )


@dataclass(frozen=True)
class CoordinateTransform:
    x_axis: AxisTransform = AxisTransform()
    y_axis: AxisTransform = AxisTransform()

    def apply_x(self, value: float) -> float:
        return self.x_axis.apply(value)

    def apply_y(self, value: float) -> float:
        return self.y_axis.apply(value)

    def compose(self, other: "CoordinateTransform") -> "CoordinateTransform":
        return CoordinateTransform(
            x_axis=self.x_axis.compose(other.x_axis),
            y_axis=self.y_axis.compose(other.y_axis),
        )


@dataclass(frozen=True)
class RawBounds:
    x: float
    y: float
    w: float
    h: float


_IDENTITY_TRANSFORM = CoordinateTransform()


def _rgb_to_hex(rgb: Any) -> str | None:
    if rgb is None:
        return None
    try:
        return f"#{int(rgb[0]):02X}{int(rgb[1]):02X}{int(rgb[2]):02X}"
    except Exception:
        text = str(rgb).strip()
        if len(text) == 6:
            return f"#{text.upper()}"
    return None


def _shape_type_name(shape: Any) -> str:
    shape_type = getattr(shape, "shape_type", None)
    return str(getattr(shape_type, "name", shape_type))


def _clamp_value(value: float, limit: float) -> float:
    return max(0.0, min(value, limit))


class TemplateExtractor:
    """PPTX 파일에서 콘텐츠 슬라이드의 elements JSON을 추출한다."""

    def extract(
        self,
        pptx_path: str,
        *,
        debug_output_path: str | None = None,
    ) -> list[dict[str, Any]]:
        prs = Presentation(pptx_path)
        selected_slides = self._select_content_slides(list(prs.slides))
        source_width = float(prs.slide_width)
        source_height = float(prs.slide_height)

        result: list[dict[str, Any]] = []
        debug_payload: dict[str, Any] = {
            "pptx_path": str(pptx_path),
            "source_slide_size": {
                "width_emu": int(prs.slide_width),
                "height_emu": int(prs.slide_height),
                "width_inches": round(Emu(prs.slide_width).inches, 2),
                "height_inches": round(Emu(prs.slide_height).inches, 2),
            },
            "content_slide_numbers": [slide_number for slide_number, _ in selected_slides],
            "slides": [],
        }

        for slide_number, slide in selected_slides:
            page_data, slide_debug = self._extract_slide(
                slide,
                slide_number=slide_number,
                source_width=source_width,
                source_height=source_height,
            )
            debug_payload["slides"].append(slide_debug)
            if page_data["elements"]:
                result.append(page_data)

        if debug_output_path:
            output_path = Path(debug_output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(
                json.dumps(debug_payload, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )

        return result

    def _select_content_slides(self, slides: list[Any]) -> list[tuple[int, Any]]:
        numbered_slides = list(enumerate(slides, start=1))
        if len(numbered_slides) >= 5:
            return numbered_slides[2:-2]
        return numbered_slides

    def _extract_slide(
        self,
        slide: Any,
        *,
        slide_number: int,
        source_width: float,
        source_height: float,
    ) -> tuple[dict[str, Any], dict[str, Any]]:
        page_data: dict[str, Any] = {
            "background": self._extract_background(slide),
            "elements": [],
        }
        slide_debug: dict[str, Any] = {
            "slide_number": slide_number,
            "background": page_data["background"],
            "shape_count": len(slide.shapes),
            "extracted": [],
            "skipped": [],
        }

        for shape_index, shape in enumerate(slide.shapes, start=1):
            source_path = f"slide[{slide_number}].shape[{shape_index}]"
            page_data["elements"].extend(
                self._extract_shape_element(
                    shape,
                    source_width=source_width,
                    source_height=source_height,
                    transform=_IDENTITY_TRANSFORM,
                    slide_debug=slide_debug,
                    source_path=source_path,
                )
            )

        slide_debug["extracted_element_count"] = len(page_data["elements"])
        if not page_data["elements"]:
            return page_data, slide_debug

        validated = PageLayout.model_validate(page_data).model_dump(
            exclude={"slots"},
            exclude_none=True,
        )
        slide_debug["validated_element_count"] = len(validated["elements"])
        _LOGGER.debug(
            "Extracted template slide %s with %s elements (%s skipped)",
            slide_number,
            len(validated["elements"]),
            len(slide_debug["skipped"]),
        )
        return validated, slide_debug

    def _extract_background(self, slide: Any) -> str:
        fill = slide.background.fill
        if fill is not None and fill.type == MSO_FILL.SOLID:
            color = _rgb_to_hex(getattr(fill.fore_color, "rgb", None))
            if color:
                return color
        return _DEFAULT_BACKGROUND

    def _extract_shape_element(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
        slide_debug: dict[str, Any],
        source_path: str,
    ) -> list[dict[str, Any]]:
        shape_type = getattr(shape, "shape_type", None)

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            return self._extract_group_shape(
                shape,
                source_width=source_width,
                source_height=source_height,
                transform=transform,
                slide_debug=slide_debug,
                source_path=source_path,
            )

        if shape_type == MSO_SHAPE_TYPE.LINE:
            line_element = self._extract_line_shape(
                shape,
                source_width=source_width,
                source_height=source_height,
                transform=transform,
            )
            if line_element is None:
                self._record_skip(slide_debug, source_path, shape, "line_out_of_bounds")
                return []
            self._record_element(slide_debug, source_path, shape, line_element, "connector")
            return [line_element]

        if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return self._extract_auto_shape(
                shape,
                source_width=source_width,
                source_height=source_height,
                transform=transform,
                slide_debug=slide_debug,
                source_path=source_path,
            )

        if shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return self._extract_freeform_shape(
                shape,
                source_width=source_width,
                source_height=source_height,
                transform=transform,
                slide_debug=slide_debug,
                source_path=source_path,
            )

        if shape_type in {MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER}:
            return self._extract_text_only_shape(
                shape,
                source_width=source_width,
                source_height=source_height,
                transform=transform,
                slide_debug=slide_debug,
                source_path=source_path,
            )

        if shape_type in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.DIAGRAM,
            MSO_SHAPE_TYPE.CANVAS,
        }:
            self._record_skip(slide_debug, source_path, shape, "unsupported_shape_type")
            return []

        self._record_skip(slide_debug, source_path, shape, "unhandled_shape_type")
        return []

    def _extract_group_shape(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
        slide_debug: dict[str, Any],
        source_path: str,
    ) -> list[dict[str, Any]]:
        group_transform = transform.compose(self._group_transform(shape))
        elements: list[dict[str, Any]] = []

        for child_index, child in enumerate(shape.shapes, start=1):
            elements.extend(
                self._extract_shape_element(
                    child,
                    source_width=source_width,
                    source_height=source_height,
                    transform=group_transform,
                    slide_debug=slide_debug,
                    source_path=f"{source_path}.child[{child_index}]",
                )
            )

        if not elements:
            self._record_skip(slide_debug, source_path, shape, "group_without_extractable_children")
        return elements

    def _extract_auto_shape(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
        slide_debug: dict[str, Any],
        source_path: str,
    ) -> list[dict[str, Any]]:
        if float(shape.width) == 0 or float(shape.height) == 0:
            line_element = self._extract_line_from_box_shape(
                shape,
                source_width=source_width,
                source_height=source_height,
                transform=transform,
            )
            if line_element is None:
                self._record_skip(slide_debug, source_path, shape, "auto_shape_line_out_of_bounds")
                return []
            self._record_element(slide_debug, source_path, shape, line_element, "auto_shape_line")
            return [line_element]

        bounds = self._extract_box_bounds(shape, source_width, source_height, transform)
        if bounds is None:
            self._record_skip(slide_debug, source_path, shape, "shape_out_of_bounds")
            return []

        elements: list[dict[str, Any]] = []
        fill_color = self._extract_fill_color(shape)
        text_payload = self._extract_text_payload(shape, bounds)

        if text_payload is None:
            if fill_color is None:
                self._record_skip(slide_debug, source_path, shape, "auto_shape_without_solid_fill_or_text")
                return []
            shape_element = self._build_shape_element(shape, fill_color, bounds)
            if shape_element is None:
                self._record_skip(slide_debug, source_path, shape, "shape_payload_build_failed")
                return []
            self._record_element(slide_debug, source_path, shape, shape_element, "auto_shape_fill")
            return [shape_element]

        if fill_color:
            shape_element = self._build_shape_element(shape, fill_color, bounds)
            if shape_element:
                elements.append(shape_element)
                self._record_element(slide_debug, source_path, shape, shape_element, "text_background")

        elements.append(text_payload)
        self._record_element(slide_debug, source_path, shape, text_payload, "text_content")
        return elements

    def _extract_text_only_shape(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
        slide_debug: dict[str, Any],
        source_path: str,
    ) -> list[dict[str, Any]]:
        bounds = self._extract_box_bounds(shape, source_width, source_height, transform)
        if bounds is None:
            self._record_skip(slide_debug, source_path, shape, "text_shape_out_of_bounds")
            return []

        text_payload = self._extract_text_payload(shape, bounds)
        if text_payload is None:
            self._record_skip(slide_debug, source_path, shape, "text_shape_without_text")
            return []

        self._record_element(slide_debug, source_path, shape, text_payload, "text_only_shape")
        return [text_payload]

    def _extract_line_shape(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
    ) -> dict[str, Any] | None:
        bounds = self._extract_line_bounds(shape, source_width, source_height, transform)
        if bounds is None:
            return None
        return {
            "type": "line",
            **bounds,
            "line_color": self._extract_line_color(shape) or _DEFAULT_LINE_COLOR,
            "line_width": self._extract_line_width(shape),
        }

    def _extract_line_from_box_shape(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
    ) -> dict[str, Any] | None:
        raw_bounds = RawBounds(
            x=float(shape.left),
            y=float(shape.top),
            w=float(shape.width),
            h=float(shape.height),
        )
        bounds = self._normalize_line_box_bounds(raw_bounds, source_width, source_height, transform)
        if bounds is None:
            return None
        return {
            "type": "line",
            **bounds,
            "line_color": self._extract_line_color(shape) or _DEFAULT_LINE_COLOR,
            "line_width": self._extract_line_width(shape),
        }

    def _extract_freeform_shape(
        self,
        shape: Any,
        *,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
        slide_debug: dict[str, Any],
        source_path: str,
    ) -> list[dict[str, Any]]:
        bounds = self._extract_box_bounds(shape, source_width, source_height, transform)
        if bounds is None:
            self._record_skip(slide_debug, source_path, shape, "freeform_out_of_bounds")
            return []

        fill_type = getattr(getattr(shape, "fill", None), "type", None)
        fill_color = self._extract_fill_color(shape)
        text_payload = self._extract_text_payload(shape, bounds)

        if fill_type == MSO_FILL.PICTURE and text_payload is None:
            self._record_skip(slide_debug, source_path, shape, "freeform_picture_fill")
            return []

        if self._looks_like_line(bounds):
            line_element = self._build_line_from_box_bounds(
                bounds,
                line_color=self._extract_line_color(shape) or fill_color or _DEFAULT_LINE_COLOR,
                line_width=self._extract_line_width(shape),
            )
            if line_element:
                self._record_element(slide_debug, source_path, shape, line_element, "freeform_as_line")
                return [line_element]

        elements: list[dict[str, Any]] = []
        if text_payload is not None:
            if fill_color:
                shape_element = self._build_shape_element(shape, fill_color, bounds)
                if shape_element:
                    elements.append(shape_element)
                    self._record_element(slide_debug, source_path, shape, shape_element, "freeform_text_background")
            elements.append(text_payload)
            self._record_element(slide_debug, source_path, shape, text_payload, "freeform_text")
            return elements

        if fill_color:
            shape_element = self._build_shape_element(shape, fill_color, bounds)
            if shape_element:
                self._record_element(slide_debug, source_path, shape, shape_element, "freeform_as_shape")
                return [shape_element]

        self._record_skip(slide_debug, source_path, shape, "freeform_without_supported_mapping")
        return []

    def _extract_text_payload(
        self,
        shape: Any,
        bounds: dict[str, float],
    ) -> dict[str, Any] | None:
        if not getattr(shape, "has_text_frame", False):
            return None

        text_frame = shape.text_frame
        paragraphs = self._text_paragraphs(text_frame)
        if not paragraphs:
            return None

        if self._is_bullet_list(paragraphs):
            return self._build_bullet_list_element(paragraphs, text_frame, bounds)
        return self._build_text_box_element("\n".join(paragraphs), text_frame, bounds)

    def _build_shape_element(
        self,
        shape: Any,
        fill_color: str,
        bounds: dict[str, float],
    ) -> dict[str, Any] | None:
        shape_type = "rectangle"
        try:
            if getattr(shape, "auto_shape_type", None) == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                shape_type = "round_rectangle"
        except ValueError:
            shape_type = "rectangle"

        return {
            "type": "shape",
            "shape_type": shape_type,
            **bounds,
            "fill_color": fill_color,
        }

    def _build_text_box_element(
        self,
        text: str,
        text_frame: Any,
        bounds: dict[str, float],
    ) -> dict[str, Any]:
        font_data = self._extract_font_data(text_frame)
        return {
            "type": "text_box",
            **bounds,
            "text": text.strip(),
            **font_data,
        }

    def _build_bullet_list_element(
        self,
        paragraphs: list[str],
        text_frame: Any,
        bounds: dict[str, float],
    ) -> dict[str, Any]:
        font_data = self._extract_font_data(text_frame)
        return {
            "type": "bullet_list",
            **bounds,
            "items": paragraphs[:5],
            "bullet_char": "-",
            "bullet_color": _DEFAULT_BULLET_COLOR,
            "font_name": font_data["font_name"],
            "font_size": font_data["font_size"],
            "font_color": font_data["font_color"],
        }

    def _build_line_from_box_bounds(
        self,
        bounds: dict[str, float],
        *,
        line_color: str,
        line_width: float,
    ) -> dict[str, Any] | None:
        if bounds["w"] >= bounds["h"]:
            y = round(bounds["y"] + (bounds["h"] / 2), 2)
            if bounds["w"] <= 0:
                return None
            return {
                "type": "line",
                "x": bounds["x"],
                "y": y,
                "w": bounds["w"],
                "h": 0.0,
                "line_color": line_color,
                "line_width": line_width,
            }

        x = round(bounds["x"] + (bounds["w"] / 2), 2)
        if bounds["h"] <= 0:
            return None
        return {
            "type": "line",
            "x": x,
            "y": bounds["y"],
            "w": 0.0,
            "h": bounds["h"],
            "line_color": line_color,
            "line_width": line_width,
        }

    def _extract_box_bounds(
        self,
        shape: Any,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
    ) -> dict[str, float] | None:
        raw_bounds = RawBounds(
            x=float(shape.left),
            y=float(shape.top),
            w=float(shape.width),
            h=float(shape.height),
        )
        return self._normalize_box_bounds(raw_bounds, source_width, source_height, transform)

    def _normalize_box_bounds(
        self,
        raw_bounds: RawBounds,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform = _IDENTITY_TRANSFORM,
    ) -> dict[str, float] | None:
        left = self._normalize_x(transform.apply_x(raw_bounds.x), source_width)
        top = self._normalize_y(transform.apply_y(raw_bounds.y), source_height)
        right = self._normalize_x(transform.apply_x(raw_bounds.x + raw_bounds.w), source_width)
        bottom = self._normalize_y(transform.apply_y(raw_bounds.y + raw_bounds.h), source_height)

        left = _clamp_value(left, SLIDE_WIDTH_INCHES)
        top = _clamp_value(top, SLIDE_HEIGHT_INCHES)
        right = _clamp_value(right, SLIDE_WIDTH_INCHES)
        bottom = _clamp_value(bottom, SLIDE_HEIGHT_INCHES)

        x = round(min(left, right), 2)
        y = round(min(top, bottom), 2)
        w = round(abs(right - left), 2)
        h = round(abs(bottom - top), 2)
        if w <= 0 or h <= 0:
            return None
        return {"x": x, "y": y, "w": w, "h": h}

    def _extract_line_bounds(
        self,
        shape: Any,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform,
    ) -> dict[str, float] | None:
        try:
            start_x = float(shape.begin_x)
            start_y = float(shape.begin_y)
            end_x = float(shape.end_x)
            end_y = float(shape.end_y)
        except Exception:
            raw_bounds = RawBounds(
                x=float(shape.left),
                y=float(shape.top),
                w=float(shape.width),
                h=float(shape.height),
            )
            return self._normalize_line_box_bounds(raw_bounds, source_width, source_height, transform)

        start_x = _clamp_value(
            self._normalize_x(transform.apply_x(start_x), source_width),
            SLIDE_WIDTH_INCHES,
        )
        start_y = _clamp_value(
            self._normalize_y(transform.apply_y(start_y), source_height),
            SLIDE_HEIGHT_INCHES,
        )
        end_x = _clamp_value(
            self._normalize_x(transform.apply_x(end_x), source_width),
            SLIDE_WIDTH_INCHES,
        )
        end_y = _clamp_value(
            self._normalize_y(transform.apply_y(end_y), source_height),
            SLIDE_HEIGHT_INCHES,
        )

        width = round(end_x - start_x, 2)
        height = round(end_y - start_y, 2)
        if width == 0 and height == 0:
            return None
        return {
            "x": round(start_x, 2),
            "y": round(start_y, 2),
            "w": width,
            "h": height,
        }

    def _normalize_line_box_bounds(
        self,
        raw_bounds: RawBounds,
        source_width: float,
        source_height: float,
        transform: CoordinateTransform = _IDENTITY_TRANSFORM,
    ) -> dict[str, float] | None:
        start_x = _clamp_value(
            self._normalize_x(transform.apply_x(raw_bounds.x), source_width),
            SLIDE_WIDTH_INCHES,
        )
        start_y = _clamp_value(
            self._normalize_y(transform.apply_y(raw_bounds.y), source_height),
            SLIDE_HEIGHT_INCHES,
        )
        end_x = _clamp_value(
            self._normalize_x(transform.apply_x(raw_bounds.x + raw_bounds.w), source_width),
            SLIDE_WIDTH_INCHES,
        )
        end_y = _clamp_value(
            self._normalize_y(transform.apply_y(raw_bounds.y + raw_bounds.h), source_height),
            SLIDE_HEIGHT_INCHES,
        )

        width = round(end_x - start_x, 2)
        height = round(end_y - start_y, 2)
        if width == 0 and height == 0:
            return None
        return {
            "x": round(start_x, 2),
            "y": round(start_y, 2),
            "w": width,
            "h": height,
        }

    def _group_transform(self, group_shape: Any) -> CoordinateTransform:
        xfrm = getattr(getattr(group_shape, "_element", None), "xfrm", None)
        if xfrm is None:
            return CoordinateTransform(
                x_axis=AxisTransform(translate=float(group_shape.left)),
                y_axis=AxisTransform(translate=float(group_shape.top)),
            )

        child_width = float(getattr(getattr(xfrm, "chExt", None), "cx", 0) or 0)
        child_height = float(getattr(getattr(xfrm, "chExt", None), "cy", 0) or 0)
        group_width = float(getattr(xfrm, "cx", 0) or 0)
        group_height = float(getattr(xfrm, "cy", 0) or 0)
        child_offset_x = float(getattr(getattr(xfrm, "chOff", None), "x", 0) or 0)
        child_offset_y = float(getattr(getattr(xfrm, "chOff", None), "y", 0) or 0)
        group_offset_x = float(getattr(xfrm, "x", 0) or 0)
        group_offset_y = float(getattr(xfrm, "y", 0) or 0)

        scale_x = (group_width / child_width) if child_width else 1.0
        scale_y = (group_height / child_height) if child_height else 1.0

        return CoordinateTransform(
            x_axis=AxisTransform(
                scale=scale_x,
                translate=group_offset_x - (child_offset_x * scale_x),
            ),
            y_axis=AxisTransform(
                scale=scale_y,
                translate=group_offset_y - (child_offset_y * scale_y),
            ),
        )

    def _extract_fill_color(self, shape: Any) -> str | None:
        fill = getattr(shape, "fill", None)
        if fill is None or fill.type != MSO_FILL.SOLID:
            return None
        return _rgb_to_hex(getattr(fill.fore_color, "rgb", None))

    def _extract_line_color(self, shape: Any) -> str | None:
        try:
            line_fill = shape.line.fill
        except Exception:
            return None
        if line_fill is None or line_fill.type != MSO_FILL.SOLID:
            return None
        return _rgb_to_hex(getattr(line_fill.fore_color, "rgb", None))

    def _extract_line_width(self, shape: Any) -> float:
        try:
            width = shape.line.width
        except Exception:
            width = None
        if width is None or int(width) <= 0:
            return _DEFAULT_LINE_WIDTH
        return round(max(_DEFAULT_LINE_WIDTH, Emu(int(width)).pt), 2)

    def _extract_font_data(self, text_frame: Any) -> dict[str, Any]:
        paragraph, run = self._find_first_text_run(text_frame)
        font = run.font if run is not None else getattr(paragraph, "font", None)

        font_name = getattr(font, "name", None) or _DEFAULT_FONT_NAME
        font_size = getattr(font, "size", None)
        font_size_value = int(round(font_size.pt)) if font_size is not None else _DEFAULT_FONT_SIZE
        font_bold = bool(getattr(font, "bold", False))
        font_color = _rgb_to_hex(getattr(getattr(font, "color", None), "rgb", None)) or _DEFAULT_FONT_COLOR
        alignment = self._extract_alignment(paragraph)

        return {
            "font_name": font_name,
            "font_size": font_size_value,
            "font_bold": font_bold,
            "font_color": font_color,
            "align": alignment,
        }

    def _find_first_text_run(self, text_frame: Any) -> tuple[Any, Any | None]:
        for paragraph in text_frame.paragraphs:
            paragraph_text = (paragraph.text or "").strip()
            if paragraph.runs:
                for run in paragraph.runs:
                    if (run.text or "").strip():
                        return paragraph, run
            if paragraph_text:
                return paragraph, None
        return text_frame.paragraphs[0], None

    def _extract_alignment(self, paragraph: Any) -> str:
        alignment = getattr(paragraph, "alignment", None)
        if alignment == PP_ALIGN.CENTER:
            return "center"
        if alignment == PP_ALIGN.RIGHT:
            return "right"
        return "left"

    def _text_paragraphs(self, text_frame: Any) -> list[str]:
        return [
            text
            for text in ((paragraph.text or "").strip() for paragraph in text_frame.paragraphs)
            if text
        ]

    def _is_bullet_list(self, paragraphs: list[str]) -> bool:
        return len(paragraphs) >= 2 and all(len(text) <= 200 for text in paragraphs)

    def _looks_like_line(self, bounds: dict[str, float]) -> bool:
        short_side = min(bounds["w"], bounds["h"])
        long_side = max(bounds["w"], bounds["h"])
        if short_side <= 0 or long_side <= 0:
            return False
        return (
            short_side <= _FREEFORM_LINE_THICKNESS_MAX
            and (long_side / short_side) >= _FREEFORM_LINE_ASPECT_RATIO
        )

    def _normalize_x(self, value: float, source_width: float) -> float:
        return (value / source_width) * SLIDE_WIDTH_INCHES if source_width else 0.0

    def _normalize_y(self, value: float, source_height: float) -> float:
        return (value / source_height) * SLIDE_HEIGHT_INCHES if source_height else 0.0

    def _record_element(
        self,
        slide_debug: dict[str, Any],
        source_path: str,
        shape: Any,
        element: dict[str, Any],
        note: str,
    ) -> None:
        slide_debug["extracted"].append(
            {
                "source": source_path,
                "shape_name": getattr(shape, "name", ""),
                "shape_type": _shape_type_name(shape),
                "element_type": element.get("type"),
                "note": note,
            }
        )

    def _record_skip(
        self,
        slide_debug: dict[str, Any],
        source_path: str,
        shape: Any,
        reason: str,
    ) -> None:
        slide_debug["skipped"].append(
            {
                "source": source_path,
                "shape_name": getattr(shape, "name", ""),
                "shape_type": _shape_type_name(shape),
                "reason": reason,
            }
        )

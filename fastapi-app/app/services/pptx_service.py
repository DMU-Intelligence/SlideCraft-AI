from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..models.project_state import ProjectState
from ..schemas.generate import (
    BulletListElement,
    ImagePlaceholderElement,
    PageLayout,
    ShapeElement,
    SlideContent,
    SlideElement,
    TextBoxElement,
)

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

_ALIGN_MAP: dict[str, Any] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

_THEMES: dict[str, dict[str, str | int]] = {
    # 청색 계열 단일 색조 — 채도·명도 변화만으로 구성
    "clean_light": {
        "background": "#F2F5FD",
        "accent": "#2563EB",
        "text": "#0D1A38",
        "muted": "#4B72C8",
        "panel": "#E4EBFB",
        "soft_panel": "#C8D8F8",
        "inverse_text": "#FFFFFF",
        "title_size": 27,
        "body_size": 18,
    },
    # 짙은 남색 계열 단일 색조 — 어두운 배경에 밝은 청색 강조
    "bold_dark": {
        "background": "#0D1528",
        "accent": "#5A8FFA",
        "text": "#DCE6FA",
        "muted": "#7EA8F0",
        "panel": "#162244",
        "soft_panel": "#1E3060",
        "inverse_text": "#FFFFFF",
        "title_size": 28,
        "body_size": 18,
    },
    # 인디고 계열 단일 색조 — 쿨한 보라·슬레이트 계열 밝기·채도 변화
    "editorial": {
        "background": "#F0F1FA",
        "accent": "#4F46E5",
        "text": "#1E1B4B",
        "muted": "#6366F1",
        "panel": "#E0E7FF",
        "soft_panel": "#C7D2FE",
        "inverse_text": "#FFFFFF",
        "title_size": 27,
        "body_size": 18,
    },
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    if len(value) != 6:
        value = "FFFFFF"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _theme_for(name: str) -> dict[str, str | int]:
    return _THEMES.get(name, _THEMES["clean_light"])


def _set_bg(slide: Any, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_text_box(
    slide: Any,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    font_color: str,
    bold: bool = False,
    align: str = "left",
    font_name: str = "Malgun Gothic",
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    if not str(text).strip():
        return
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = text_box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = _ALIGN_MAP.get(align, PP_ALIGN.LEFT)
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _hex_to_rgb(font_color)


def _add_panel(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    radius: bool = False,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()


def _add_title_box(slide: Any, title: str, theme: dict[str, str | int], *, top: float = 0.72) -> None:
    _add_panel(slide, left=0.75, top=top, width=4.1, height=0.62, fill_color=str(theme["accent"]))
    _add_text_box(
        slide,
        text=title,
        left=1.0,
        top=top + 0.12,
        width=3.55,
        height=0.32,
        font_size=19,
        font_color=str(theme["inverse_text"]),
        bold=True,
    )


def _render_text_box(slide: Any, elem: TextBoxElement) -> None:
    _add_text_box(
        slide,
        text=elem.text,
        left=elem.x,
        top=elem.y,
        width=elem.w,
        height=elem.h,
        font_size=elem.font_size,
        font_color=elem.font_color,
        bold=elem.font_bold,
        align=elem.align,
        font_name=elem.font_name,
    )


def _render_shape(slide: Any, elem: ShapeElement) -> None:
    shape_map = {
        "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "round_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    }
    shape = slide.shapes.add_shape(
        shape_map.get(elem.shape_type, MSO_AUTO_SHAPE_TYPE.RECTANGLE),
        Inches(elem.x),
        Inches(elem.y),
        Inches(elem.w),
        Inches(elem.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(elem.fill_color)
    shape.line.fill.background()


def _render_bullet_list(slide: Any, elem: BulletListElement) -> None:
    text_box = slide.shapes.add_textbox(Inches(elem.x), Inches(elem.y), Inches(elem.w), Inches(elem.h))
    frame = text_box.text_frame
    frame.word_wrap = True
    bullet_rgb = _hex_to_rgb(elem.bullet_color)
    font_rgb = _hex_to_rgb(elem.font_color)
    bullet_prefix = f"{elem.bullet_char or '-'} "

    for index, item_text in enumerate(elem.items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(4)
        paragraph.space_after = Pt(3)

        bullet = paragraph.add_run()
        bullet.text = bullet_prefix
        bullet.font.name = elem.font_name
        bullet.font.size = Pt(elem.font_size)
        bullet.font.color.rgb = bullet_rgb
        bullet.font.bold = True

        body = paragraph.add_run()
        body.text = item_text
        body.font.name = elem.font_name
        body.font.size = Pt(elem.font_size)
        body.font.color.rgb = font_rgb


def _render_image_placeholder(slide: Any, elem: ImagePlaceholderElement) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(elem.x),
        Inches(elem.y),
        Inches(elem.w),
        Inches(elem.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb("#E4EBFB")
    shape.line.color.rgb = _hex_to_rgb("#7EA8F0")
    shape.line.width = Pt(1.2)

    lines = [f"[ {elem.label} ]"]
    if elem.description:
        lines.append(elem.description)

    label_top = elem.y + elem.h / 2 - 0.35
    _add_text_box(
        slide,
        text="\n".join(lines),
        left=elem.x + 0.1,
        top=label_top,
        width=elem.w - 0.2,
        height=0.75,
        font_size=13,
        font_color="#4B72C8",
        align="center",
    )


def _render_element(slide: Any, element: SlideElement) -> None:
    if isinstance(element, TextBoxElement):
        _render_text_box(slide, element)
    elif isinstance(element, ShapeElement):
        _render_shape(slide, element)
    elif isinstance(element, BulletListElement):
        _render_bullet_list(slide, element)
    elif isinstance(element, ImagePlaceholderElement):
        _render_image_placeholder(slide, element)


def _coerce_list(value: Any) -> list[str]:
    if not isinstance(value, list):
        return []
    return [str(item).strip() for item in value if str(item).strip()]


def _render_bullets(
    slide: Any,
    items: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: dict[str, str | int],
) -> None:
    if not items:
        return
    _render_bullet_list(
        slide,
        BulletListElement(
            x=left,
            y=top,
            w=width,
            h=height,
            items=items[:5],
            bullet_char="•",
            bullet_color=str(theme["accent"]),
            font_color=str(theme["text"]),
            font_size=int(theme["body_size"]),
        ),
    )


def _render_people(
    slide: Any,
    people: list[str],
    *,
    left: float,
    top: float,
    width: float,
    theme: dict[str, str | int],
    align: str = "left",
) -> None:
    if not people:
        return
    _add_text_box(
        slide,
        text="\n".join(people),
        left=left,
        top=top,
        width=width,
        height=1.2,
        font_size=15,
        font_color=str(theme["muted"]),
        align=align,
    )


def _render_title_page(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    # Vertical accent bar
    _add_panel(slide, left=0.7, top=0.7, width=0.18, height=5.9, fill_color=str(theme["accent"]))
    # Main headline — no eyebrow label
    _add_text_box(
        slide,
        text=str(slots.get("headline", "")),
        left=1.1,
        top=1.55,
        width=8.8,
        height=1.55,
        font_size=int(theme["title_size"]) + 6,
        font_color=str(theme["text"]),
        bold=True,
    )
    _add_text_box(
        slide,
        text=str(slots.get("body", "")),
        left=1.1,
        top=3.35,
        width=7.8,
        height=0.7,
        font_size=int(theme["body_size"]),
        font_color=str(theme["muted"]),
    )
    people = _coerce_list(slots.get("people"))
    if people:
        _add_text_box(
            slide,
            text="\n".join(people),
            left=9.2,
            top=5.8,
            width=3.0,
            height=0.65,
            font_size=11,
            font_color=str(theme["muted"]),
            align="right",
        )


def _render_content_box_list(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    # Title label box (top-left accent bar)
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    # Body subtitle (right of title)
    _add_text_box(
        slide,
        text=str(slots.get("body", "")),
        left=5.1,
        top=0.87,
        width=6.7,
        height=0.45,
        font_size=14,
        font_color=str(theme["muted"]),
    )
    highlight = str(slots.get("highlight", "")).strip()
    bullets_top: float
    if highlight:
        # Full-width accent banner with highlight text
        _add_panel(slide, left=0.75, top=1.65, width=11.83, height=1.1, fill_color=str(theme["accent"]))
        _add_text_box(
            slide,
            text=highlight,
            left=1.05,
            top=1.75,
            width=11.23,
            height=0.82,
            font_size=int(theme["body_size"]),
            font_color=str(theme["inverse_text"]),
            bold=True,
        )
        bullets_top = 3.0
    else:
        bullets_top = 1.75
    # Thin vertical accent stripe beside bullet list
    _add_panel(
        slide,
        left=0.75,
        top=bullets_top,
        width=0.06,
        height=7.0 - bullets_top,
        fill_color=str(theme["accent"]),
    )
    # Bullet list on clean background — no white panel behind
    _render_bullets(
        slide,
        _coerce_list(slots.get("bullets")),
        left=1.1,
        top=bullets_top + 0.1,
        width=11.23,
        height=7.0 - bullets_top - 0.3,
        theme=theme,
    )


def _render_content_two_panel(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    _add_panel(slide, left=0.9, top=1.85, width=5.45, height=4.55, fill_color=str(theme["panel"]))
    _add_panel(slide, left=6.75, top=1.85, width=5.45, height=4.55, fill_color=str(theme["soft_panel"]))
    _render_bullets(
        slide,
        _coerce_list(slots.get("left_points") or slots.get("bullets")),
        left=1.2,
        top=2.2,
        width=4.85,
        height=3.7,
        theme=theme,
    )
    _render_bullets(
        slide,
        _coerce_list(slots.get("right_points")),
        left=7.05,
        top=2.2,
        width=4.85,
        height=3.7,
        theme=theme,
    )


def _render_content_sidebar(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    _add_panel(slide, left=0.9, top=1.7, width=2.55, height=4.95, fill_color=str(theme["accent"]))
    _add_text_box(
        slide,
        text=str(slots.get("highlight", "")) or str(slots.get("body", "")),
        left=1.2,
        top=2.0,
        width=1.95,
        height=4.2,
        font_size=16,
        font_color=str(theme["inverse_text"]),
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_panel(slide, left=3.75, top=1.7, width=8.45, height=4.95, fill_color=str(theme["panel"]))
    _render_bullets(
        slide,
        _coerce_list(slots.get("bullets")),
        left=4.1,
        top=2.05,
        width=7.75,
        height=3.9,
        theme=theme,
    )


def _render_content_split_band(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    _add_panel(slide, left=0.9, top=1.85, width=11.2, height=1.15, fill_color=str(theme["soft_panel"]))
    _add_text_box(
        slide,
        text=str(slots.get("body", "")),
        left=1.2,
        top=2.18,
        width=10.6,
        height=0.45,
        font_size=15,
        font_color=str(theme["text"]),
        bold=True,
        align="center",
    )
    _add_panel(slide, left=0.9, top=3.35, width=5.35, height=3.15, fill_color=str(theme["panel"]))
    _add_panel(slide, left=6.75, top=3.35, width=5.35, height=3.15, fill_color=str(theme["panel"]))
    _render_bullets(
        slide,
        _coerce_list(slots.get("left_points") or slots.get("bullets")),
        left=1.2,
        top=3.7,
        width=4.75,
        height=2.4,
        theme=theme,
    )
    _render_bullets(
        slide,
        _coerce_list(slots.get("right_points")),
        left=7.05,
        top=3.7,
        width=4.75,
        height=2.4,
        theme=theme,
    )


def _render_content_compact(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    _add_panel(slide, left=0.9, top=1.8, width=11.2, height=4.85, fill_color=str(theme["panel"]))
    _render_bullets(
        slide,
        _coerce_list(slots.get("bullets")),
        left=1.2,
        top=2.15,
        width=6.4,
        height=3.9,
        theme=theme,
    )
    _add_panel(slide, left=8.2, top=2.15, width=3.25, height=3.75, fill_color=str(theme["soft_panel"]))
    _add_text_box(
        slide,
        text=str(slots.get("highlight", "")) or str(slots.get("body", "")),
        left=8.5,
        top=2.55,
        width=2.65,
        height=2.9,
        font_size=15,
        font_color=str(theme["text"]),
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def _render_content_card_grid(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    """불릿 개수(2~5)에 따라 개별 카드를 자동 배치한다."""
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    body = str(slots.get("body", "")).strip()
    if body:
        _add_text_box(
            slide, text=body, left=5.1, top=0.87, width=6.7, height=0.45,
            font_size=14, font_color=str(theme["muted"]),
        )

    items = _coerce_list(slots.get("bullets") or slots.get("left_points"))
    if not items:
        return
    items = items[:5]
    count = len(items)

    MARGIN_L, MARGIN_T = 0.75, 1.55
    AVAIL_W, AVAIL_H = 11.83, 5.7
    GAP = 0.28

    def _card_bounds(n: int) -> list[tuple[float, float, float, float]]:
        if n <= 3:
            cw = (AVAIL_W - (n - 1) * GAP) / n
            ch = AVAIL_H
            return [(MARGIN_L + i * (cw + GAP), MARGIN_T, cw, ch) for i in range(n)]
        if n == 4:
            cw = (AVAIL_W - GAP) / 2
            ch = (AVAIL_H - GAP) / 2
            return [
                (MARGIN_L + (i % 2) * (cw + GAP), MARGIN_T + (i // 2) * (ch + GAP), cw, ch)
                for i in range(4)
            ]
        # 5: 3 on top + 2 on bottom centered
        cw3 = (AVAIL_W - 2 * GAP) / 3
        cw2 = (AVAIL_W - GAP) / 2
        ch = (AVAIL_H - GAP) / 2
        offset2 = (AVAIL_W - (2 * cw2 + GAP)) / 2
        top_cards = [(MARGIN_L + i * (cw3 + GAP), MARGIN_T, cw3, ch) for i in range(3)]
        bot_cards = [(MARGIN_L + offset2 + i * (cw2 + GAP), MARGIN_T + ch + GAP, cw2, ch) for i in range(2)]
        return top_cards + bot_cards

    for i, (cx, cy, cw, ch) in enumerate(_card_bounds(count)):
        _add_panel(slide, left=cx, top=cy, width=cw, height=ch, fill_color=str(theme["panel"]))
        _add_panel(slide, left=cx, top=cy, width=cw, height=0.2, fill_color=str(theme["accent"]))
        _add_text_box(
            slide, text=str(i + 1),
            left=cx + 0.18, top=cy + 0.28, width=0.6, height=0.45,
            font_size=22, font_color=str(theme["accent"]), bold=True,
        )
        _add_text_box(
            slide, text=items[i],
            left=cx + 0.18, top=cy + 0.78, width=cw - 0.35, height=ch - 0.95,
            font_size=int(theme["body_size"]) - 1, font_color=str(theme["text"]),
        )


def _render_content_steps(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    """순서가 있는 단계를 좌측 번호 박스 + 우측 텍스트 패널로 배치한다."""
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    body = str(slots.get("body", "")).strip()
    if body:
        _add_text_box(
            slide, text=body, left=5.1, top=0.87, width=6.7, height=0.45,
            font_size=14, font_color=str(theme["muted"]),
        )

    items = _coerce_list(slots.get("bullets") or slots.get("left_points"))
    if not items:
        return
    items = items[:5]
    count = len(items)

    MARGIN_L, MARGIN_T = 0.75, 1.6
    AVAIL_H = 5.6
    MIN_GAP = 0.22
    MAX_STEP_H = 1.35
    NUM_BOX_W = 0.85

    raw_h = (AVAIL_H - (count - 1) * MIN_GAP) / count
    step_h = min(MAX_STEP_H, raw_h)
    total_used = count * step_h + (count - 1) * MIN_GAP
    extra = AVAIL_H - total_used
    gap = MIN_GAP + extra / max(1, count - 1)

    text_panel_w = 11.83 - NUM_BOX_W - 0.22

    for i, item_text in enumerate(items):
        cy = MARGIN_T + i * (step_h + gap)
        # Number box
        _add_panel(slide, left=MARGIN_L, top=cy, width=NUM_BOX_W, height=step_h,
                   fill_color=str(theme["accent"]))
        _add_text_box(
            slide, text=str(i + 1),
            left=MARGIN_L, top=cy + (step_h - 0.48) / 2,
            width=NUM_BOX_W, height=0.48,
            font_size=24, font_color=str(theme["inverse_text"]),
            bold=True, align="center",
        )
        # Connector line (except last step)
        if i < count - 1:
            line_x = MARGIN_L + NUM_BOX_W / 2 - 0.04
            _add_panel(slide, left=line_x, top=cy + step_h, width=0.08, height=gap,
                       fill_color=str(theme["soft_panel"]))
        # Text panel
        px = MARGIN_L + NUM_BOX_W + 0.22
        _add_panel(slide, left=px, top=cy, width=text_panel_w, height=step_h,
                   fill_color=str(theme["panel"]))
        _add_text_box(
            slide, text=item_text,
            left=px + 0.28, top=cy + (step_h - 0.5) / 2,
            width=text_panel_w - 0.4, height=step_h - 0.1,
            font_size=int(theme["body_size"]) - 1, font_color=str(theme["text"]),
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )


def _render_content_highlight_split(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    """좌측 accent 강조 패널 + 우측 불릿 콘텐츠 분할 레이아웃."""
    _add_title_box(slide, str(slots.get("title_box_label") or slots.get("headline", "")), theme)
    # Left accent panel
    LEFT_W = 3.9
    _add_panel(slide, left=0.75, top=1.55, width=LEFT_W, height=5.7, fill_color=str(theme["accent"]))
    highlight = str(slots.get("highlight", "") or slots.get("body", "")).strip()
    if highlight:
        _add_text_box(
            slide, text=highlight,
            left=0.98, top=2.3, width=LEFT_W - 0.45, height=3.8,
            font_size=int(theme["body_size"]) + 2, font_color=str(theme["inverse_text"]),
            bold=True, vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
    # Right content panel
    RIGHT_X = 0.75 + LEFT_W + 0.3
    RIGHT_W = 11.83 - LEFT_W - 0.3
    _add_panel(slide, left=RIGHT_X, top=1.55, width=RIGHT_W, height=5.7,
               fill_color=str(theme["panel"]))
    _render_bullets(
        slide,
        _coerce_list(slots.get("bullets")),
        left=RIGHT_X + 0.3, top=1.95, width=RIGHT_W - 0.45, height=4.9,
        theme=theme,
    )


def _render_closing_page(slide: Any, slots: dict[str, Any], theme: dict[str, str | int]) -> None:
    _add_panel(slide, left=0.9, top=1.0, width=11.5, height=5.2, fill_color=str(theme["panel"]))
    _add_panel(slide, left=0.9, top=1.0, width=11.5, height=0.2, fill_color=str(theme["accent"]))
    _add_text_box(
        slide,
        text=str(slots.get("headline", "Thank you")),
        left=1.35,
        top=2.0,
        width=10.4,
        height=0.9,
        font_size=int(theme["title_size"]) + 10,
        font_color=str(theme["text"]),
        bold=True,
        align="center",
    )
    _add_text_box(
        slide,
        text=str(slots.get("body", "")),
        left=2.0,
        top=3.15,
        width=9.1,
        height=0.5,
        font_size=16,
        font_color=str(theme["muted"]),
        align="center",
    )
    _render_people(
        slide,
        _coerce_list(slots.get("people")),
        left=2.2,
        top=4.2,
        width=8.8,
        theme=theme,
        align="center",
    )


def _render_slots(slide: Any, slide_content: SlideContent, page: PageLayout) -> None:
    theme = _theme_for(slide_content.theme)
    slots = page.slots
    if not slots:
        return

    _set_bg(slide, page.background or str(theme["background"]))

    variant = slide_content.slide_variant
    if variant == "title_page":
        _render_title_page(slide, slots, theme)
    elif variant == "content_box_list":
        _render_content_box_list(slide, slots, theme)
    elif variant == "content_two_panel":
        _render_content_two_panel(slide, slots, theme)
    elif variant == "content_sidebar":
        _render_content_sidebar(slide, slots, theme)
    elif variant == "content_split_band":
        _render_content_split_band(slide, slots, theme)
    elif variant == "content_compact":
        _render_content_compact(slide, slots, theme)
    elif variant == "content_card_grid":
        _render_content_card_grid(slide, slots, theme)
    elif variant == "content_steps":
        _render_content_steps(slide, slots, theme)
    elif variant == "content_highlight_split":
        _render_content_highlight_split(slide, slots, theme)
    elif variant == "closing_page":
        _render_closing_page(slide, slots, theme)


def _build_page(prs: Presentation, slide_content: SlideContent, page: PageLayout) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = _theme_for(slide_content.theme)
    _set_bg(slide, page.background or str(theme["background"]))

    if page.elements:
        for element in page.elements:
            _render_element(slide, element)
    elif page.slots:
        _render_slots(slide, slide_content, page)


class PptxGenerator:
    def generate(self, state: ProjectState) -> bytes:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        for slide_content in state.slides:
            for page in slide_content.pages:
                _build_page(prs, slide_content, page)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

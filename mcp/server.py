import os
import io
import uuid
import asyncio
from typing import Dict, List, Optional
from mcp.server import Server
from mcp.types import Tool, TextContent
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Create the MCP server
app = Server("pptx-generator")

# Store for active presentations (in memory state management)
presentations: Dict[str, Presentation] = {}

# Helper: Hex to RGBColor
def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "FFFFFF"  # fallback
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)

@app.list_tools()
async def handle_list_tools() -> list[Tool]:
    """도구 목록 정의: 순수 pptx 라이브러리 기능만 제공"""
    return [
        Tool(
            name="get_full_presentation_json",
            description="프레젠테이션의 모든 슬라이드와 각 슬라이드 내 요소(텍스트, 위치, 크기 등)를 JSON 형식으로 반환합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"}
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="get_presentation_info",
            description="프레젠테이션의 슬라이드 개수와 레이아웃 목록을 가져옵니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"}
                },
                "required": ["presentation_id"]
            }
        ),
        
        Tool(
            name="create_presentation",
            description="새로운 PPT 프레젠테이션 객체를 메모리에 생성하거나 템플릿을 불러오고 presentation_id를 반환합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "width_inches": {"type": "number", "default": 13.33},
                    "height_inches": {"type": "number", "default": 7.5},
                    "template_path": {"type": "string", "description": "불러올 템플릿 파일(.pptx)의 경로 (선택 사항)"}
                }
            }
        ),
        Tool(
            name="add_slide",
            description="프레젠테이션에 슬라이드를 추가합니다. 반환된 slide_index를 기억해야 합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "layout_index": {"type": "integer", "default": 6}
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="set_slide_background",
            description="특정 슬라이드의 배경색을 설정합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "hex_color": {"type": "string", "description": "예: #FFFFFF"}
                },
                "required": ["presentation_id", "slide_index", "hex_color"]
            }
        ),
        Tool(
            name="add_text_box",
            description="슬라이드에 텍스트 상자를 추가합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "left": {"type": "number"},
                    "top": {"type": "number"},
                    "width": {"type": "number"},
                    "height": {"type": "number"},
                    "text": {"type": "string"},
                    "font_size": {"type": "integer", "default": 24},
                    "font_name": {"type": "string", "default": "Arial"},
                    "font_color": {"type": "string", "default": "#000000"},
                    "font_bold": {"type": "boolean", "default": False},
                    "align": {"type": "string", "enum": ["left", "center", "right"], "default": "left"}
                },
                "required": ["presentation_id", "slide_index", "left", "top", "width", "height", "text"]
            }
        ),
        Tool(
            name="add_shape",
            description="슬라이드에 사각형 도형을 추가합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "left": {"type": "number"},
                    "top": {"type": "number"},
                    "width": {"type": "number"},
                    "height": {"type": "number"},
                    "fill_color": {"type": "string", "default": "#FFFFFF"}
                },
                "required": ["presentation_id", "slide_index", "left", "top", "width", "height"]
            }
        ),
        Tool(
            name="add_bullet_list",
            description="슬라이드에 글머리 기호(불릿 리스트) 텍스트 상자를 추가합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "left": {"type": "number"},
                    "top": {"type": "number"},
                    "width": {"type": "number"},
                    "height": {"type": "number"},
                    "items": {"type": "array", "items": {"type": "string"}},
                    "font_size": {"type": "integer", "default": 18},
                    "font_name": {"type": "string", "default": "Arial"},
                    "font_color": {"type": "string", "default": "#000000"},
                    "bullet_color": {"type": "string", "default": "#000000"},
                    "bullet_char": {"type": "string", "default": "•"}
                },
                "required": ["presentation_id", "slide_index", "left", "top", "width", "height", "items"]
            }
        ),
        Tool(
            name="save_presentation",
            description="프레젠테이션을 디스크에 .pptx 파일로 저장합니다.",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "filename": {"type": "string"}
                },
                "required": ["presentation_id", "filename"]
            }
        )
    ]

@app.call_tool()
async def handle_call_tool(name: str, arguments: dict | None) -> list[TextContent]:
    """도구 실행: 순수 pptx 조작만 수행"""
    if not arguments:
        arguments = {}

    if name == "create_presentation":
        template_path = arguments.get("template_path")
        prs_id = str(uuid.uuid4())
        
        if template_path:
            if not os.path.exists(template_path):
                return [TextContent(type="text", text=f"Error: Template file not found at {template_path}")]
            try:
                prs = Presentation(template_path)
            except Exception as e:
                return [TextContent(type="text", text=f"Error loading template: {str(e)}")]
        else:
            w = arguments.get("width_inches", 13.33)
            h = arguments.get("height_inches", 7.5)
            prs = Presentation()
            prs.slide_width = Inches(w)
            prs.slide_height = Inches(h)
            
        presentations[prs_id] = prs
        return [TextContent(type="text", text=f"{prs_id}")]

    elif name == "get_presentation_info":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        
        layouts = [{"index": i, "name": layout.name} for i, layout in enumerate(prs.slide_layouts)]
        slide_count = len(prs.slides)
        
        import json
        info = {
            "slide_count": slide_count,
            "layouts": layouts
        }
        return [TextContent(type="text", text=json.dumps(info, indent=2, ensure_ascii=False))]

    elif name == "get_full_presentation_json":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        
        def get_color_info(color_obj):
            try:
                if not color_obj: return None
                res = {"type": str(color_obj.type)}
                if hasattr(color_obj, "rgb") and color_obj.rgb:
                    res["rgb"] = f"#{str(color_obj.rgb)}"
                if hasattr(color_obj, "theme_color") and color_obj.theme_color:
                    res["theme_color"] = str(color_obj.theme_color)
                return res
            except:
                return None

        def get_fill_info(fill_obj):
            try:
                res = {"type": str(fill_obj.type)}
                if hasattr(fill_obj, "fore_color"):
                    res["fore_color"] = get_color_info(fill_obj.fore_color)
                return res
            except:
                return None

        slides_data = []
        for i, slide in enumerate(prs.slides):
            bg_info = None
            try:
                bg_info = get_fill_info(slide.background.fill)
            except: pass

            elements = []
            for shape in slide.shapes:
                el = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "left": shape.left.inches if shape.left else 0,
                    "top": shape.top.inches if shape.top else 0,
                    "width": shape.width.inches if shape.width else 0,
                    "height": shape.height.inches if shape.height else 0,
                }

                try:
                    if hasattr(shape, "fill"):
                        el["fill"] = get_fill_info(shape.fill)
                    if hasattr(shape, "line"):
                        el["line"] = {
                            "color": get_color_info(shape.line.color),
                            "width_pt": shape.line.width.pt if shape.line.width else 0,
                            "dash_style": str(shape.line.dash_style) if shape.line.dash_style else None
                        }
                except: pass

                if shape.has_text_frame:
                    el["text"] = shape.text
                    paragraphs = []
                    for p in shape.text_frame.paragraphs:
                        p_data = {
                            "alignment": str(p.alignment),
                            "runs": []
                        }
                        for run in p.runs:
                            r_data = {
                                "text": run.text,
                                "font": {
                                    "name": run.font.name,
                                    "size_pt": run.font.size.pt if run.font.size else None,
                                    "bold": run.font.bold,
                                    "italic": run.font.italic,
                                    "color": get_color_info(run.font.color)
                                }
                            }
                            p_data["runs"].append(r_data)
                        paragraphs.append(p_data)
                    el["paragraphs"] = paragraphs
                
                elements.append(el)
            
            slides_data.append({
                "slide_index": i,
                "layout_name": slide.slide_layout.name,
                "background": bg_info,
                "elements": elements
            })
            
        import json
        return [TextContent(type="text", text=json.dumps(slides_data, indent=2, ensure_ascii=False))]

    elif name == "add_slide":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        layout_idx = arguments.get("layout_index", 6)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide_index = len(prs.slides) - 1
        return [TextContent(type="text", text=f"{slide_index}")]

    elif name == "set_slide_background":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        slide_idx = arguments.get("slide_index")
        slide = prs.slides[slide_idx]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(arguments.get("hex_color"))
        return [TextContent(type="text", text="Background set successfully.")]

    elif name == "add_text_box":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        slide_idx = arguments.get("slide_index")
        slide = prs.slides[slide_idx]
        
        txBox = slide.shapes.add_textbox(
            Inches(arguments.get("left")), Inches(arguments.get("top")), 
            Inches(arguments.get("width")), Inches(arguments.get("height"))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(arguments.get("align", "left"), PP_ALIGN.LEFT)
        
        run = p.add_run()
        run.text = arguments.get("text")
        run.font.name = arguments.get("font_name", "Arial")
        run.font.size = Pt(arguments.get("font_size", 24))
        run.font.bold = arguments.get("font_bold", False)
        run.font.color.rgb = _hex_to_rgb(arguments.get("font_color", "#000000"))
        return [TextContent(type="text", text="Text box added successfully.")]

    elif name == "add_shape":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        slide_idx = arguments.get("slide_index")
        slide = prs.slides[slide_idx]
        
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE_TYPE.RECTANGLE
            Inches(arguments.get("left")), Inches(arguments.get("top")), 
            Inches(arguments.get("width")), Inches(arguments.get("height"))
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(arguments.get("fill_color", "#FFFFFF"))
        shape.line.fill.background()
        return [TextContent(type="text", text="Shape added successfully.")]

    elif name == "add_bullet_list":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        slide_idx = arguments.get("slide_index")
        slide = prs.slides[slide_idx]
        
        txBox = slide.shapes.add_textbox(
            Inches(arguments.get("left")), Inches(arguments.get("top")), 
            Inches(arguments.get("width")), Inches(arguments.get("height"))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        
        bullet_rgb = _hex_to_rgb(arguments.get("bullet_color", "#000000"))
        font_rgb = _hex_to_rgb(arguments.get("font_color", "#000000"))
        font_name = arguments.get("font_name", "Arial")
        font_size = arguments.get("font_size", 18)
        bullet_char = arguments.get("bullet_char", "•")

        for i, item_text in enumerate(arguments.get("items", [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(4)

            dot = p.add_run()
            dot.text = f"{bullet_char}  "
            dot.font.name = font_name
            dot.font.size = Pt(font_size)
            dot.font.color.rgb = bullet_rgb
            dot.font.bold = True

            body = p.add_run()
            body.text = item_text
            body.font.name = font_name
            body.font.size = Pt(font_size)
            body.font.color.rgb = font_rgb
            body.font.bold = False
            
        return [TextContent(type="text", text="Bullet list added successfully.")]

    elif name == "save_presentation":
        prs_id = arguments.get("presentation_id")
        if prs_id not in presentations: return [TextContent(type="text", text="Error: Presentation not found.")]
        prs = presentations[prs_id]
        
        filename = arguments.get("filename")
        if not filename.endswith(".pptx"): filename += ".pptx"
        
        os.makedirs(os.path.join(os.getcwd(), "output"), exist_ok=True)
        save_path = os.path.join(os.getcwd(), "output", filename)
        prs.save(save_path)
        
        # 메모리 정리 (선택적)
        del presentations[prs_id]
        
        return [TextContent(type="text", text=f"Presentation saved to: {save_path}")]

    raise ValueError(f"Unknown tool: {name}")

if __name__ == "__main__":
    from mcp.server.stdio import stdio_server
    
    async def main():
        async with stdio_server() as (read_stream, write_stream):
            await app.run(read_stream, write_stream, app.create_initialization_options())
            
    asyncio.run(main())

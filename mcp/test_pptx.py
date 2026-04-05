from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "FFFFFF"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)

def test_gen():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb("#F0F0F0")
    
    # Text Box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "MCP PPTX Generator Test"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb("#333333")
    
    # Bullet List
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(3))
    tf2 = txBox2.text_frame
    items = ["Feature 1: MCP Integration", "Feature 2: Dynamic Generation", "Feature 3: python-pptx support"]
    
    for i, item_text in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run = p.add_run()
        run.text = f"•  {item_text}"
        run.font.size = Pt(24)
    
    os.makedirs("output", exist_ok=True)
    prs.save("output/test_output.pptx")
    print("Test PPTX generated at output/test_output.pptx")

if __name__ == "__main__":
    test_gen()
